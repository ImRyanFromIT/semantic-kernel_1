from pathlib import Path

from pptx import Presentation

from semantic-kernel_poc.app.services.extractors.pptx_extractor import extract_org_chart_snippets


def make_simple_org_chart(tmp_path: Path) -> Path:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Data Storage Services"
    txbox = slide.shapes.add_textbox(100, 150, 500, 300)
    tf = txbox.text_frame
    tf.text = "Backup Services\nMission:\nSpecialist team managing backup solutions\nServices:\nBackup Design & Implementation\nTechnologies:\nVeeam, CommVault"
    path = tmp_path / "orgchart.pptx"
    prs.save(path)
    return path


def test_extract_org_chart_snippets(tmp_path: Path):
    pptx_path = make_simple_org_chart(tmp_path)
    snippets, teams = extract_org_chart_snippets(pptx_path)
    assert any(s.section_type == "mission" for s in snippets)
    assert any(s.section_type == "services" for s in snippets)
    assert any(s.section_type == "technologies" for s in snippets)
    assert "Backup Services" in teams


