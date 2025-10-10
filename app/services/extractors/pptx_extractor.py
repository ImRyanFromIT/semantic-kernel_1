from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Dict, List, Tuple

from pptx import Presentation

from ..models.entities import DocumentSnippet


HEADING_ALIASES: Dict[str, str] = {
    "mission": "mission",
    "scope": "mission",
    "mandate": "mission",
    "services": "services",
    "offerings": "services",
    "technologies": "technologies",
}

LEVEL_RE = re.compile(r"^(?P<name>[^()]+?)\s*\((?P<level>L\d)\)$", re.IGNORECASE)


def _normalize_heading(text: str) -> str | None:
    key = text.strip().lower().rstrip(":")
    return HEADING_ALIASES.get(key)


def extract_org_chart_snippets(pptx_path: str | Path) -> Tuple[List[DocumentSnippet], List[str]]:
    """Extract text sections from an org chart deck into snippets.

    Returns a list of snippets and a list of inferred team names. The caller can
    associate snippets with teams by deck structure or an external mapping.
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    prs = Presentation(str(path))
    snippets: List[DocumentSnippet] = []
    inferred_teams: List[str] = []

    for slide in prs.slides:
        slide_title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        # Naive heuristic: columns/panels per team; collect textbox content
        panel_idx = 0
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            panel_idx += 1
            raw_text = shape.text_frame.text
            lines = [ln.strip() for ln in raw_text.splitlines() if ln.strip()]
            if not lines:
                continue

            team_name = lines[0]
            if team_name not in inferred_teams:
                inferred_teams.append(team_name)

            # Detect sections by heading markers in subsequent lines
            current_section: str | None = None
            buffer: List[str] = []

            def flush():  # type: ignore[no-redef]
                nonlocal buffer, current_section
                if current_section and buffer:
                    snippets.append(
                        DocumentSnippet(
                            source=str(path.name),
                            title=f"{team_name} - {current_section}",
                            content=" ".join(buffer),
                            team_id=team_name.lower().replace(" ", "-"),
                            section_type=current_section,
                        )
                    )
                buffer = []

            for ln in lines[1:]:
                norm = _normalize_heading(ln)
                if norm:
                    flush()
                    current_section = norm
                    continue
                buffer.append(ln)

            flush()

    logging.getLogger(__name__).info(
        "Extracted %d snippets and %d team names from PPTX %s", len(snippets), len(inferred_teams), path
    )
    return snippets, inferred_teams


